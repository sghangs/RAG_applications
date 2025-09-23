# ingestion/parsers/pptx_parser.py
from pptx import Presentation
import io
from typing import List, Dict, Any

from unstructured.partition.pptx import partition_pptx
from unstructured.partition.image import partition_image
from utils.logger import logger
from exceptions import KnowledgeManagementException


class HybridPPTXParser:
    """
    Hybrid PPTX parser:
    - Extracts text from shapes in slides
    - Extracts embedded images and applies OCR on each image individually
    - Falls back to full slide-deck OCR if overall text is below threshold
    """

    def __init__(self, text_threshold: int = 50, enable_ocr: bool = True):
        self.text_threshold = text_threshold
        self.enable_ocr = enable_ocr

    # --------------------------
    # Native text extraction
    # --------------------------
    def _extract_slide_text(self, slide) -> str:
        """Extract visible text from all shapes in a slide."""
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        return "\n".join(texts).strip()

    # --------------------------
    # Image extraction
    # --------------------------
    def _extract_image_blobs(self, slide) -> List[Dict[str, Any]]:
        """
        Extract image blobs from a slide.
        Returns: [{"image_index": int, "blob": bytes}]
        """
        image_blobs: List[Dict[str, Any]] = []
        idx = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                idx += 1
                try:
                    blob = shape.image.blob
                    image_blobs.append({
                        "image_index": idx,
                        "blob": blob,
                    })
                except Exception as e:
                    logger.debug(f"Skipping unreadable PPTX image: {e}")
        return image_blobs

    # --------------------------
    # OCR logic
    # --------------------------
    def _ocr_images_only(self, image_blobs: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Run OCR on each image blob individually using partition_image."""
        ocr_blocks: List[Dict[str, Any]] = []
        if not image_blobs:
            return ocr_blocks

        for img in image_blobs:
            img_idx = img.get("image_index")
            blob = img.get("blob")
            if not blob:
                continue

            try:
                elements = partition_image(
                    file_content=blob,
                    strategy="hi_res",
                    ocr_strategy="ocr_only"
                )
                for el in elements:
                    cat = getattr(el, "category", "text").lower()
                    text = str(el).strip()
                    if not text:
                        continue

                    if cat in ("table", "tabular"):
                        block_type = "table"
                    else:
                        block_type = "image_text"

                    ocr_blocks.append({
                        "type": block_type,
                        "text": text,
                        "metadata": {
                            "source": "pptx-image-ocr",
                            "image_index": img_idx,
                        }
                    })
            except Exception as e:
                logger.warning(f"OCR failed for PPTX image {img_idx}: {e}")
                continue

        return ocr_blocks

    def _ocr_full_deck(self, content: bytes) -> List[Dict[str, Any]]:
        """Fallback: run Unstructured OCR on the entire PPTX deck."""
        try:
            elements = partition_pptx(file=io.BytesIO(content), strategy="hi_res")
            results: List[Dict[str, Any]] = []
            for el in elements:
                cat = getattr(el, "category", "text").lower()
                text = str(el).strip()
                if not text:
                    continue
                if cat in ("table", "tabular"):
                    btype = "table"
                else:
                    btype = "text"
                results.append({
                    "type": btype,
                    "text": text,
                    "metadata": {"source": "unstructured-ocr"}
                })
            return results
        except Exception as e:
            raise KnowledgeManagementException(
                f"Full OCR extraction failed on PPTX: {e}",
                None,
                "HybridPPTXParser"
            )

    # --------------------------
    # Main entry
    # --------------------------
    def parse(self, content: bytes) -> List[Dict[str, Any]]:
        try:
            prs = Presentation(io.BytesIO(content))
            results: List[Dict[str, Any]] = []

            total_text_len = 0

            for slide_idx, slide in enumerate(prs.slides, start=1):
                # Native text
                slide_text = self._extract_slide_text(slide)
                if slide_text:
                    results.append({
                        "type": "text",
                        "text": slide_text,
                        "metadata": {"slide": slide_idx, "source": "pptx-native"}
                    })
                    total_text_len += len(slide_text)

                # Image OCR
                image_blobs = self._extract_image_blobs(slide)
                if self.enable_ocr and image_blobs:
                    ocr_blocks = self._ocr_images_only(image_blobs)
                    results.extend(ocr_blocks)
                elif image_blobs:
                    # placeholders if OCR disabled
                    for img in image_blobs:
                        results.append({
                            "type": "image",
                            "text": "",
                            "metadata": {
                                "source": "pptx-image",
                                "image_index": img.get("image_index"),
                                "slide": slide_idx,
                            }
                        })

            # Fallback to full OCR if deck text is too low
            if total_text_len < self.text_threshold and self.enable_ocr:
                logger.info("Low text in PPTX deck — running full OCR fallback")
                full_ocr_blocks = self._ocr_full_deck(content)
                existing_texts = {b["text"] for b in results if b.get("text")}
                for fb in full_ocr_blocks:
                    if fb["text"] not in existing_texts:
                        results.append(fb)

            return results

        except Exception as e:
            raise KnowledgeManagementException(
                f"Failed to parse PPTX: {e}",
                None,
                "HybridPPTXParser"
            )
