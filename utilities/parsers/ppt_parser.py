import io
from pptx import Presentation
from unstructured.partition.pptx import partition_pptx
from .base import BaseParser
from .utils_image import ocr_image_bytes

class HybridPPTXParser(BaseParser):
    def __init__(self, enable_ocr_images: bool = False):
        self.enable_ocr_images = enable_ocr_images

    def _parse_with_unstructured(self, content: bytes) -> list[dict]:
        elements = partition_pptx(file=io.BytesIO(content))
        blocks = []
        for idx, el in enumerate(elements, start=1):
            text = str(el).strip()
            if text:
                blocks.append({"text": text, "metadata": {"source": "unstructured-pptx", "block": idx}})
        return blocks

    def _parse_with_pptx(self, content: bytes) -> list[dict]:
        prs = Presentation(io.BytesIO(content))
        blocks = []
        for idx, slide in enumerate(prs.slides, start=1):
            text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            if text.strip():
                blocks.append({"text": text, "metadata": {"source": "python-pptx", "slide": idx}})
            if self.enable_ocr_images:
                for shp_idx, shape in enumerate(slide.shapes, start=1):
                    if shape.shape_type == 13:  # picture
                        img_bytes = shape.image.blob
                        ocr_blocks = ocr_image_bytes(img_bytes, "pptx-image", {"slide": idx, "shape": shp_idx})
                        blocks.extend(ocr_blocks)
        return blocks

    def parse(self, content: bytes) -> list[dict]:
        try:
            blocks = self._parse_with_unstructured(content)
            if blocks:
                return blocks
        except Exception:
            pass
        return self._parse_with_pptx(content)
